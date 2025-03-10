#!/usr/bin/env python
# -*- coding: utf-8 -*-

import os
import sys
from PyQt5.QtWidgets import (QApplication, QMainWindow, QPushButton, QLabel,
                             QVBoxLayout, QHBoxLayout, QFileDialog, QTextEdit,
                             QProgressBar, QWidget, QMessageBox, QSplitter,
                             QFrame, QSizePolicy)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QUrl, QSize
from PyQt5.QtGui import QFont, QIcon, QDragEnterEvent, QDropEvent, QPalette, QColor
import markdown
from pptx import Presentation

class PPTXConverter(QThread):
    progress_signal = pyqtSignal(int)
    finished_signal = pyqtSignal(str)
    error_signal = pyqtSignal(str)
    
    def __init__(self, pptx_path, output_path=None):
        super().__init__()
        self.pptx_path = pptx_path
        self.output_path = output_path or os.path.splitext(pptx_path)[0] + '.md'
        
    def run(self):
        try:
            # 打开PPT文件
            prs = Presentation(self.pptx_path)
            total_slides = len(prs.slides)
            
            # 创建Markdown内容
            md_content = f"# {os.path.basename(self.pptx_path)}\n\n"
            
            # 遍历每一张幻灯片
            for i, slide in enumerate(prs.slides):
                # 更新进度
                self.progress_signal.emit(int((i / total_slides) * 100))
                
                # 添加幻灯片标题
                if slide.shapes.title:
                    md_content += f"## {slide.shapes.title.text}\n\n"
                else:
                    md_content += f"## Slide {i+1}\n\n"
                
                # 提取幻灯片内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip() and shape != slide.shapes.title:
                        md_content += f"{shape.text}\n\n"
                
                # 添加分隔符
                md_content += "---\n\n"
            
            # 保存Markdown文件
            with open(self.output_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            self.progress_signal.emit(100)
            self.finished_signal.emit(self.output_path)
            
        except Exception as e:
            self.error_signal.emit(str(e))

class DropArea(QWidget):
    file_dropped = pyqtSignal(str)
    
    def __init__(self):
        super().__init__()
        self.setAcceptDrops(True)
        layout = QVBoxLayout()
        self.label = QLabel("拖放PPTX文件到这里")
        self.label.setAlignment(Qt.AlignCenter)
        font = QFont()
        font.setPointSize(18)
        self.label.setFont(font)
        layout.addWidget(self.label)
        self.setLayout(layout)
        self.setMinimumHeight(150)  # 减小高度以留出更多空间给预览区域
        self.setStyleSheet("""
            border: 2px dashed #aaa; 
            border-radius: 8px;
            padding: 20px;
            background-color: #f8f9fa;
        """)
    
    def dragEnterEvent(self, event: QDragEnterEvent):
        if event.mimeData().hasUrls() and event.mimeData().urls()[0].toLocalFile().endswith('.pptx'):
            event.acceptProposedAction()
            self.setStyleSheet("""
                border: 2px dashed #3498db; 
                border-radius: 8px; 
                padding: 20px;
                background-color: #e8f4fc;
            """)
    
    def dragLeaveEvent(self, event):
        self.setStyleSheet("""
            border: 2px dashed #aaa; 
            border-radius: 8px;
            padding: 20px;
            background-color: #f8f9fa;
        """)
    
    def dropEvent(self, event: QDropEvent):
        file_path = event.mimeData().urls()[0].toLocalFile()
        self.file_dropped.emit(file_path)
        self.setStyleSheet("""
            border: 2px dashed #aaa; 
            border-radius: 8px;
            padding: 20px;
            background-color: #f8f9fa;
        """)

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.initUI()
        self.current_file = None
        self.converter = None
        self.setStyleSheet("""
            QMainWindow, QWidget {
                background-color: #f5f5f7;
                color: #333333;
            }
            QLabel {
                font-size: 18px;
                font-weight: bold;
                color: #2c3e50;
                padding: 5px 0;
            }
            QPushButton {
                background-color: #3498db;
                color: white;
                border: none;
                border-radius: 4px;
                padding: 10px 18px;
                font-weight: bold;
                min-width: 120px;
                font-size: 16px;
            }
            QPushButton:hover {
                background-color: #2980b9;
            }
            QPushButton:disabled {
                background-color: #bdc3c7;
            }
            QProgressBar {
                border: 1px solid #bdc3c7;
                border-radius: 3px;
                text-align: center;
                height: 20px;
            }
            QProgressBar::chunk {
                background-color: #3498db;
                width: 10px;
                margin: 0.5px;
            }
            QTextEdit {
                border: 1px solid #dcdde1;
                border-radius: 4px;
                background-color: white;
                font-family: "Consolas", "Courier New", monospace;
                font-size: 16px;
                padding: 10px;
            }
            QSplitter::handle {
                background-color: #bdc3c7;
                height: 4px;
            }
        """)
    
    def initUI(self):
        self.setWindowTitle('PPTX转Markdown工具')
        self.setMinimumSize(900, 700)  # 增加窗口初始大小
        
        # 创建主窗口部件和布局
        main_widget = QWidget()
        main_layout = QVBoxLayout(main_widget)
        main_layout.setSpacing(10)  # 增加布局间距
        main_layout.setContentsMargins(15, 15, 15, 15)  # 增加边距
        
        # 创建拖放区域
        self.drop_area = DropArea()
        self.drop_area.file_dropped.connect(self.process_file)
        main_layout.addWidget(self.drop_area)
        
        # 创建按钮区域
        button_layout = QHBoxLayout()
        button_layout.setSpacing(10)  # 按钮之间的间距
        self.select_btn = QPushButton('选择PPTX文件')
        self.select_btn.setIcon(QIcon.fromTheme('document-open'))
        self.select_btn.clicked.connect(self.select_file)
        self.select_btn.setCursor(Qt.PointingHandCursor)  # 鼠标悬停时显示手型光标
        
        self.convert_btn = QPushButton('转换为Markdown')
        self.convert_btn.setIcon(QIcon.fromTheme('document-save'))
        self.convert_btn.clicked.connect(self.convert_file)
        self.convert_btn.setEnabled(False)
        self.convert_btn.setCursor(Qt.PointingHandCursor)  # 鼠标悬停时显示手型光标
        
        button_layout.addWidget(self.select_btn)
        button_layout.addWidget(self.convert_btn)
        button_layout.addStretch(1)  # 添加弹性空间，使按钮靠左对齐
        main_layout.addLayout(button_layout)
        
        # 创建进度条
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.progress_bar.setTextVisible(True)  # 显示进度百分比
        main_layout.addWidget(self.progress_bar)
        
        # 添加一个水平分隔线
        line = QFrame()
        line.setFrameShape(QFrame.HLine)
        line.setFrameShadow(QFrame.Sunken)
        line.setStyleSheet("background-color: #bdc3c7;")
        main_layout.addWidget(line)
        
        # 创建分割器
        splitter = QSplitter(Qt.Vertical)
        splitter.setHandleWidth(8)  # 增加分割条宽度，便于拖动
        splitter.setChildrenCollapsible(False)  # 防止拖动时完全折叠子部件
        
        # 创建预览区域
        preview_widget = QWidget()
        preview_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)  # 允许预览区域扩展
        preview_layout = QVBoxLayout(preview_widget)
        preview_layout.setContentsMargins(0, 10, 0, 0)  # 调整边距
        
        preview_header = QHBoxLayout()
        preview_label = QLabel('Markdown预览')
        preview_header.addWidget(preview_label)
        
        # 添加缩放按钮
        zoom_in_md = QPushButton("+")
        zoom_in_md.setFixedSize(28, 28)
        zoom_in_md.setToolTip("放大字体")
        zoom_in_md.clicked.connect(lambda: self.zoom_text(self.preview_text, 1))
        
        zoom_out_md = QPushButton("-")
        zoom_out_md.setFixedSize(28, 28)
        zoom_out_md.setToolTip("缩小字体")
        zoom_out_md.clicked.connect(lambda: self.zoom_text(self.preview_text, -1))
        
        preview_header.addStretch(1)
        preview_header.addWidget(zoom_out_md)
        preview_header.addWidget(zoom_in_md)
        preview_layout.addLayout(preview_header)
        
        self.preview_text = QTextEdit()
        self.preview_text.setReadOnly(True)
        self.preview_text.setMinimumHeight(250)  # 增加最小高度
        self.preview_text.setStyleSheet("""
            font-family: "Consolas", "Courier New", monospace;
            font-size: 18px;
            line-height: 1.8;
            padding: 15px;
        """)
        preview_layout.addWidget(self.preview_text)
        splitter.addWidget(preview_widget)
        
        # 创建HTML预览区域
        html_widget = QWidget()
        html_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)  # 允许预览区域扩展
        html_layout = QVBoxLayout(html_widget)
        html_layout.setContentsMargins(0, 10, 0, 0)  # 调整边距
        
        html_header = QHBoxLayout()
        html_label = QLabel('HTML预览')
        html_header.addWidget(html_label)
        
        # 添加缩放按钮
        zoom_in_html = QPushButton("+")
        zoom_in_html.setFixedSize(28, 28)
        zoom_in_html.setToolTip("放大内容")
        zoom_in_html.clicked.connect(lambda: self.zoom_text(self.html_preview, 1))
        
        zoom_out_html = QPushButton("-")
        zoom_out_html.setFixedSize(28, 28)
        zoom_out_html.setToolTip("缩小内容")
        zoom_out_html.clicked.connect(lambda: self.zoom_text(self.html_preview, -1))
        
        html_header.addStretch(1)
        html_header.addWidget(zoom_out_html)
        html_header.addWidget(zoom_in_html)
        html_layout.addLayout(html_header)
        
        self.html_preview = QTextEdit()
        self.html_preview.setReadOnly(True)
        self.html_preview.setMinimumHeight(250)  # 增加最小高度
        html_layout.addWidget(self.html_preview)
        splitter.addWidget(html_widget)
        
        # 设置分割器初始大小
        splitter.setSizes([350, 350])  # 增加初始分配空间
        
        main_layout.addWidget(splitter, 1)  #
    
        # 设置中央部件
        self.setCentralWidget(main_widget)
    
    def select_file(self):
        file_path, _ = QFileDialog.getOpenFileName(
            self, '选择PPTX文件', '', 'PowerPoint Files (*.pptx)'
        )
        if file_path:
            self.process_file(file_path)
    
    def process_file(self, file_path):
        self.current_file = file_path
        self.convert_btn.setEnabled(True)
        self.statusBar().showMessage(f'已选择文件: {os.path.basename(file_path)}')
        self.drop_area.label.setText(f'已选择: {os.path.basename(file_path)}')
    
    def convert_file(self):
        if not self.current_file:
            return
        
        output_path, _ = QFileDialog.getSaveFileName(
            self, '保存Markdown文件', 
            os.path.splitext(self.current_file)[0] + '.md',
            'Markdown Files (*.md)'
        )
        
        if not output_path:
            return
        
        # 设置UI状态
        self.progress_bar.setValue(0)
        self.progress_bar.setVisible(True)
        self.convert_btn.setEnabled(False)
        self.select_btn.setEnabled(False)
        self.statusBar().showMessage('正在转换...')
        
        # 创建并启动转换线程
        self.converter = PPTXConverter(self.current_file, output_path)
        self.converter.progress_signal.connect(self.update_progress)
        self.converter.finished_signal.connect(self.conversion_finished)
        self.converter.error_signal.connect(self.conversion_error)
        self.converter.start()
    
    def update_progress(self, value):
        self.progress_bar.setValue(value)
    
    def conversion_finished(self, output_path):
        # 恢复UI状态
        self.progress_bar.setVisible(False)
        self.convert_btn.setEnabled(True)
        self.select_btn.setEnabled(True)
        self.statusBar().showMessage(f'转换完成: {output_path}')
        
        # 显示预览
        try:
            with open(output_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # 设置Markdown预览，使用更大的字体和更好的格式
            self.preview_text.setStyleSheet("""
                font-family: "Consolas", "Courier New", monospace;
                font-size: 18px;
                line-height: 1.8;
                padding: 15px;
            """)
            self.preview_text.setText(md_content)
            
            # 设置HTML预览，添加基本样式
            html_content = markdown.markdown(md_content)
            styled_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <style>
                    body {{ 
                        font-family: 'Segoe UI', Arial, sans-serif; 
                        line-height: 1.8;
                        padding: 25px;
                        color: #333;
                        max-width: 100%;
                        font-size: 16px;
                    }}
                    h1 {{ font-size: 32px; color: #2c3e50; margin-bottom: 20px; }}
                    h2 {{ font-size: 28px; color: #3498db; margin-top: 25px; margin-bottom: 15px; }}
                    p {{ margin-bottom: 15px; font-size: 16px; }}
                    hr {{ border: none; border-top: 1px solid #eee; margin: 25px 0; }}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
            self.html_preview.setHtml(styled_html)
            
            # 显示成功消息，使用更友好的格式
            QMessageBox.information(self, '转换成功', 
                                   f'PPTX文件已成功转换为Markdown格式并保存到:\n{output_path}')
        except Exception as e:
            QMessageBox.warning(self, '预览错误', f'无法预览文件: {str(e)}')
    
    def zoom_text(self, text_edit, direction):
        """放大或缩小文本编辑器的字体"""
        current_font = text_edit.font()
        size = current_font.pointSize()
        new_size = max(10, size + direction * 3)  # 限制最小字体大小为10，增大缩放步长
        current_font.setPointSize(new_size)
        text_edit.setFont(current_font)
    
    def conversion_error(self, error_msg):
        # 恢复UI状态
        self.progress_bar.setVisible(False)
        self.convert_btn.setEnabled(True)
        self.select_btn.setEnabled(True)
        self.statusBar().showMessage('转换失败')
        
        # 显示错误消息
        QMessageBox.critical(self, '转换错误', f'转换过程中发生错误:\n{error_msg}')

def main():
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())

if __name__ == '__main__':
    main()